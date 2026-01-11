import openpyxl
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import pandas as pd

def sync_excel_to_ppt(excel_file: str, ppt_file: str, excel_sheet: str, add_missing_tools: bool = True):
    """
    Sync data from Excel to PowerPoint
    
    Excel columns: Tool, Service\nUse Case, Requestor, Status
    PPT columns: AI Tool, Tool Description, Requestor, Current State
    
    Params:
        excel_file: Path to Excel file
        ppt_file: Path to PowerPoint file  
        excel_sheet: Excel sheet name
        add_missing_tools: Whether to add new rows for tools not found in PPT
    """
    
    # === STYLE CONSTANTS ===
    FONT_NAME = "Calibri"
    FONT_SIZE = Pt(11)          # deck uses Calibri 11
    HEADER_BOLD = True          # set False if your headers are not bold
    BODY_BOLD = False

    # === PAGINATION CONSTANTS ===
    # Max number of *body* rows (excluding header) allowed per slide table before creating a new slide.
    MAX_BODY_ROWS = 8  # <-- tweak this to match your template's table height/spacing

    def set_cell_style(cell, *, is_header=False):
        """Apply consistent font styling to a cell's text. Call AFTER assigning cell.text."""
        tf = cell.text_frame
        if len(tf.paragraphs) == 0:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = ""
        for p in tf.paragraphs:
            if len(p.runs) == 0:
                p.add_run()
            for run in p.runs:
                run.font.name = FONT_NAME
                run.font.size = FONT_SIZE
                run.font.bold = HEADER_BOLD if is_header else BODY_BOLD
                # Optional color:
                # run.font.color.rgb = RGBColor(0, 0, 0)

    # excel col -> PPT col (kept for reference; logic matches headers directly)
    column_mapping = {
        'Tool': 'AI Tool',
        'Service\nUse Case': 'Tool Description', 
        'Requestor': 'Requestor',
        'Status': 'Current State'
    }

    print("reading xlsx")
    workbook = openpyxl.load_workbook(excel_file) # open excel file
    worksheet = workbook[excel_sheet] # open sheet from file
    
    # get data from Excel
    excel_data = []
    headers = []
    
    for row_idx, row in enumerate(worksheet.iter_rows(values_only=True)):
        if row_idx == 0:
            headers = [str(cell) if cell else "" for cell in row]
        else:
            if any(cell is not None for cell in row):  # skip empty rows
                excel_data.append([str(cell) if cell else "" for cell in row])
    
    # pandas dataFrame from the worksheet data in excel_data
    excel_df = pd.DataFrame(excel_data, columns=headers)

    # close excel file -- use dataframe from now
    workbook.close() 
    
    # check required columns in Excel
    required_excel_cols = ['Tool', 'Service\n Use Case' if 'Service\n Use Case' in excel_df.columns else 'Service\nUse Case', 'Requestor', 'Status']
    # Normalize potential header variation with/without space after \n
    excel_df.rename(columns={'Service\n Use Case': 'Service\nUse Case'}, inplace=True)
    missing_cols = [col for col in ['Tool', 'Service\nUse Case', 'Requestor', 'Status'] if col not in excel_df.columns]
    if missing_cols: # error checking
        print(f"Error: Missing Excel columns: {missing_cols}")
        return
    
    print(f"Found {len(excel_df)} tools in Excel")
    
    # open ppt
    print("opening ppt")
    presentation = Presentation(ppt_file)
    
    updates_made = 0 # track number of updates made
    rows_added = 0 # track number of rows added
    not_found_before_adding = 0 # track tools not found prior to appending
    target_table = None  # store reference to the main table for adding rows
    target_slide_idx = None
    target_shape = None     # keep the shape so we can rebuild
    target_slide = None     # keep the slide object too (needed to reinsert)
    
    required_ppt_cols = ['AI Tool', 'Tool Description', 'Requestor', 'Current State']

    def _table_headers_and_indices(table):
        """Return (headers_list, indices_map) or (None, None) if not a valid table."""
        try:
            num_rows = len(table.rows)
            num_cols = len(table.columns)
        except Exception:
            return None, None
        if num_rows == 0 or num_cols == 0:
            return None, None

        hdrs = [table.cell(0, c).text.strip() for c in range(num_cols)]
        if not all(col in hdrs for col in required_ppt_cols):
            return None, None
        idx_map = {col: hdrs.index(col) for col in required_ppt_cols}
        return hdrs, idx_map

    def _append_row_by_rebuilding(slide, table_shape, values_by_header):
        """
        Rebuild the table with +1 row and append a row with `values_by_header` dict keyed by current table headers.
        Preserves text and column widths; returns the new table (new shape.table).
        Also applies Calibri 11 styling to all cells.
        """
        table = table_shape.table
        rows = len(table.rows)
        cols = len(table.columns)

        # Capture current text grid
        grid = [[table.cell(r, c).text for c in range(cols)] for r in range(rows)]
        # Capture header order to map values correctly
        hdrs = [table.cell(0, c).text.strip() for c in range(cols)]

        # Add new row (as list of strings)
        new_row = []
        for c in range(cols):
            h = hdrs[c]
            new_row.append(values_by_header.get(h, ""))  
        grid.append(new_row)

        # Capture geometry and widths
        left, top, width, height = table_shape.left, table_shape.top, table_shape.width, table_shape.height
        col_widths = []
        try:
            for c in range(cols):
                col_widths.append(table.columns[c].width)
        except Exception:
            col_widths = None

        # Remove old table shape
        el = table_shape._element
        parent = el.getparent()
        parent.remove(el)

        # Add new table with +1 row
        new_shape = slide.shapes.add_table(rows + 1, cols, left, top, width, height)
        new_table = new_shape.table

        # Restore column widths if we have them
        if col_widths:
            for c in range(cols):
                try:
                    new_table.columns[c].width = col_widths[c]
                except Exception:
                    pass

        # Restore texts + apply styles (Calibri 11; header bold if configured)
        for r in range(rows + 1):
            for c in range(cols):
                new_table.cell(r, c).text = grid[r][c]
                set_cell_style(new_table.cell(r, c), is_header=(r == 0))

        return new_shape.table  # return the table object for continued use

    def _start_new_slide_with_table(pres, src_slide, src_shape):
        """
        Create a new slide (same layout as src_slide) and add a fresh table with the same columns/headers/widths.
        Returns (new_slide, new_shape, new_table).
        """
        # New slide with the same layout
        layout = src_slide.slide_layout
        new_slide = pres.slides.add_slide(layout)

        src_table = src_shape.table
        cols = len(src_table.columns)

        # Headers and column widths from source
        hdrs = [src_table.cell(0, c).text.strip() for c in range(cols)]
        col_widths = []
        try:
            for c in range(cols):
                col_widths.append(src_table.columns[c].width)
        except Exception:
            col_widths = None

        # Place the table at the same position/size
        left, top, width, height = src_shape.left, src_shape.top, src_shape.width, src_shape.height

        # New table with just header row for now
        new_shape = new_slide.shapes.add_table(1, cols, left, top, width, height)
        new_table = new_shape.table

        # Set headers and widths
        for c in range(cols):
            new_table.cell(0, c).text = hdrs[c]
            set_cell_style(new_table.cell(0, c), is_header=True)
            if col_widths:
                try:
                    new_table.columns[c].width = col_widths[c]
                except Exception:
                    pass

        return new_slide, new_shape, new_table

    def _ensure_capacity_or_paginate():
        """
        If the current target_table has reached MAX_BODY_ROWS, start a new slide with a fresh table,
        and update target_* references. Returns None, uses outer-scope variables via nonlocal.
        """
        nonlocal target_table, target_shape, target_slide, target_slide_idx
        if target_table is None:
            return
        body_rows = len(target_table.rows) - 1  # excluding header
        if body_rows >= MAX_BODY_ROWS:
            print(f"  Table reached {body_rows} body rows; starting a new slide for overflow…")
            new_slide, new_shape, new_table = _start_new_slide_with_table(presentation, target_slide, target_shape)
            target_table = new_table
            target_shape = new_shape
            target_slide = new_slide
            target_slide_idx = len(presentation.slides) - 1
            print(f"  Created new slide #{target_slide_idx + 1} with fresh table")

    # process each tool from Excel
    for idx, excel_row in excel_df.iterrows():
        tool_name = (excel_row['Tool'] or "").strip()
        if not tool_name:
            continue
            
        print(f"searching for tool: '{tool_name}'") # debugging here
        
        # now search for this tool in all tables in ppt 
        tool_found = False
        
        for slide_idx, slide in enumerate(presentation.slides): # iterate thru the slides
            for shape in slide.shapes: # iterate thru slide objects (shapes)
                if not getattr(shape, "has_table", False): # check if a table exists (defensive)
                    continue
                table = shape.table

                headers, col_indices = _table_headers_and_indices(table)
                if headers is None:
                    continue  # not the right table

                # store reference to the first qualifying table for adding new rows later
                if target_table is None:
                    target_table = table
                    target_slide_idx = slide_idx
                    target_shape = shape
                    target_slide = slide
                    print(f"Using table on slide {slide_idx + 1} for adding new rows")

                    # Ensure header styling is applied (in case template varies)
                    for c in range(len(table.columns)):
                        set_cell_style(table.cell(0, c), is_header=True)

                ai_tool_col_idx = col_indices.get('AI Tool')
                if ai_tool_col_idx is None:
                    continue
                    
                # check each row for the tool
                for row_idx in range(1, len(table.rows)):  # skip header row
                    current_tool = table.cell(row_idx, ai_tool_col_idx).text.strip()
                    if current_tool.lower() == tool_name.lower():
                        print(f"  Found '{tool_name}' on slide {slide_idx + 1}, row {row_idx + 1}")
                        tool_found = True
                        
                        # update the row with excel_data
                        new_desc = (excel_row['Service\nUse Case'] or "").strip()
                        desc_cell = table.cell(row_idx, col_indices['Tool Description'])
                        current_desc = desc_cell.text.strip()
                        if current_desc != new_desc:
                            desc_cell.text = new_desc
                            set_cell_style(desc_cell, is_header=False)
                            print(f"    Updated Tool Description: '{current_desc}' → '{new_desc}'")
                            updates_made += 1
                        
                        new_requester = (excel_row['Requestor'] or "").strip()
                        req_cell = table.cell(row_idx, col_indices['Requestor'])
                        current_requester = req_cell.text.strip()
                        if current_requester != new_requester:
                            req_cell.text = new_requester
                            set_cell_style(req_cell, is_header=False)
                            print(f"    Updated Requestor: '{current_requester}' → '{new_requester}'")
                            updates_made += 1
                        
                        new_status = (excel_row['Status'] or "").strip()
                        status_cell = table.cell(row_idx, col_indices['Current State'])
                        current_status = status_cell.text.strip()
                        if current_status != new_status:
                            status_cell.text = new_status
                            set_cell_style(status_cell, is_header=False)
                            print(f"    Updated Current State: '{current_status}' → '{new_status}'")
                            updates_made += 1
                        
                        # Also normalize the AI Tool cell's style so the entire row matches
                        tool_cell = table.cell(row_idx, ai_tool_col_idx)
                        set_cell_style(tool_cell, is_header=False)

                        break  # stop scanning rows in this table
                if tool_found:
                    break  # stop scanning shapes in this slide
            if tool_found:
                break  # stop scanning slides
        
        # Add missing tool immediately (paginate if needed, then append by rebuild)
        if not tool_found:
            not_found_before_adding += 1
            print(f"  Tool '{tool_name}' not found in any PowerPoint table")
            if add_missing_tools and (target_table is not None) and (target_shape is not None) and (target_slide is not None):
                # If current table is full, create a new slide/table before appending
                _ensure_capacity_or_paginate()

                # Build a dict for values keyed by table header text (so we map correctly)
                tgt_headers = [target_table.cell(0, c).text.strip() for c in range(len(target_table.columns))]
                values_by_header = {h: "" for h in tgt_headers}

                # Fill the known columns using your mapping
                col_values_from_excel = {
                    'AI Tool': tool_name,
                    'Tool Description': (excel_row['Service\nUse Case'] or "").strip(),
                    'Requestor': (excel_row['Requestor'] or "").strip(),
                    'Current State': (excel_row['Status'] or "").strip()
                }
                for h in values_by_header.keys():
                    if h in col_values_from_excel:
                        values_by_header[h] = col_values_from_excel[h]

                # Rebuild with +1 row and get back a fresh table reference (styles applied inside)
                target_table = _append_row_by_rebuilding(
                    slide=target_slide,
                    table_shape=target_shape,
                    values_by_header=values_by_header
                )

                # After rebuild, the shape object changed; refresh cached shape reference
                refreshed_shape = None
                for shp in target_slide.shapes:
                    if getattr(shp, "has_table", False):
                        hdrs, _idx = _table_headers_and_indices(shp.table)
                        if hdrs is not None:
                            refreshed_shape = shp
                            break
                if refreshed_shape is not None:
                    target_shape = refreshed_shape  # cache for subsequent appends

                rows_added += 1
                # Count as 4 updates (one per filled column) to mirror original accounting
                updates_made += 4  
            else:
                if target_table is None:
                    print("  No qualifying target table found to append new rows.")
    
    # save ppt if updates made
    if updates_made > 0:
        presentation.save(ppt_file)
        print(f"\nPowerPoint saved with {updates_made} total updates")
    else:
        print("\nNo updates were needed")
    
    print(f"\nSummary:")
    print(f"  Total tools processed: {len(excel_df)}")
    print(f"  Updates made: {updates_made}")
    print(f"  Rows added: {rows_added}")
    print(f"  Tools not found (before adding): {not_found_before_adding}")

# EDIT THIS
if __name__ == "__main__":
    excel_file = "excel.xlsx"      
    ppt_file = "ppt.pptx"  
    sheet_name = "Procurement AI tracker"                   
    
    print("starting sync...")

    sync_excel_to_ppt(excel_file, ppt_file, sheet_name)
    
    print("sync done!")
